"""
Generate CA_Office_Suite_Employee_Presentation.pptx from slide content.
Run: python docs/build_presentation.py
Requires: pip install python-pptx
"""

from __future__ import annotations

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError as exc:
    raise SystemExit("Install python-pptx: pip install python-pptx") from exc

SLIDES: list[tuple[str, list[str]]] = [
    (
        "CA Office Suite\nBusiness Rules & Logic",
        [
            "Training for office staff",
            "How the system applies rules",
        ],
    ),
    (
        "What we cover",
        [
            "Indian Financial Year (FY)",
            "Client Master & approval",
            "MIS, Tender & dashboard totals",
            "Director Mapping",
            "DIR-3 eKYC cadence",
            "Reports & permissions",
            "User accounts & branch access",
        ],
    ),
    (
        "Four main areas",
        [
            "Masters — Clients, Groups, Director Mapping, Users",
            "MIS — Fees, Receipts, Expenses, Tender, Bulk upload",
            "DIR-3 eKYC — Director KYC filings",
            "Reports — Master, MIS, Mapping, DIR-3 (export optional)",
            "Menu items depend on your permissions",
        ],
    ),
    (
        "Indian Financial Year",
        [
            "FY runs 1 April → 31 March",
            "Example: 1 Apr 2026 – 31 Mar 2027 = FY 2026-27",
            "Used for dashboard MIS, MIS Report, DIR-3 rules",
        ],
    ),
    (
        "Client Master — Approval",
        [
            "Pending — not usable in MIS, Mapping, or DIR-3",
            "Approved — full use everywhere",
            "New client by staff → pending until approver accepts",
            "Edit approved client → pending again",
            "Superuser → approved immediately",
        ],
    ),
    (
        "Client Master — Key rules",
        [
            "Director + DIN: Individual / Foreign Citizen only",
            "PAN, name, and type must be consistent",
            "Auto Client ID (e.g. AN0001)",
            "Delete blocked if MIS, mapping, or DIR-3 linked",
            "Bulk import: preview; may be pending until approved",
        ],
    ),
    (
        "Director Mapping",
        [
            "Director: approved Individual/Foreign Citizen with DIN",
            "Company: Pvt/Public Ltd, LLP, Nidhi, FPO, Sec 8",
            "One active appointment per director–company",
            "Cessation needs reason; logical dates required",
        ],
    ),
    (
        "MIS — Who & what",
        [
            "Approved clients only (from suggestions)",
            "Branch-restricted users: their branch only",
            "Fees + GST = Total (automatic)",
            "No negatives; no GST when fees = 0",
            "Tender: fees and/or deposit required",
        ],
    ),
    (
        "MIS — Dashboard FY card",
        [
            "Period: 1 April (current FY) → today",
            "Approved clients only",
            "Totals: fees (incl. GST), receipts, expenses",
            "Tender not in dashboard card — use MIS Report",
            "Same rule as Reports → MIS Report",
        ],
    ),
    (
        "MIS — Bulk upload",
        [
            "CLIENT_ID approved; name must match",
            "No GST if fees zero/blank",
            "Preview → confirm (all-or-nothing)",
            "Permissions control fees/receipts/expenses adds",
        ],
    ),
    (
        "DIR-3 eKYC — Cadence",
        [
            "Based on FY of last filing date",
            "Filing in FY 2026-27 → next from 1 April 2029",
            "System blocks earlier dates",
        ],
    ),
    (
        "DIR-3 — Dashboard pending",
        [
            "Never filed, OR",
            "Today on/after next allowed date from last filing",
            "Approved directors with DIN and director flag",
        ],
    ),
    (
        "Reports — Overview",
        [
            "MIS: approved clients only",
            "Branch filter by user access",
            "CSV export needs export permission",
            "MIS Report: multiple views from FY 2026-27",
            "Director report: active as on a date",
        ],
    ),
    (
        "Permissions & branches",
        [
            "Superuser: full access",
            "Others: per assigned permissions (access groups)",
            "Branch blank = all; else Trivandrum or Nagercoil",
            "Activity log: superusers only",
        ],
    ),
    (
        "User accounts",
        [
            "Inactive users cannot sign in",
            "Cannot deactivate/delete self or superusers",
            "First-login password change may apply",
            "Client-user email must match Client Master",
        ],
    ),
    (
        "Quick troubleshooting",
        [
            "Client not in MIS → pending or wrong branch",
            "Totals mismatch → pending excluded; FY YTD",
            "DIR-3 rejected → before next allowed date",
            "Can't delete client → still linked elsewhere",
            "Missing menu → no permission",
        ],
    ),
    (
        "Questions?",
        [
            "Full guide: CA_Office_Suite_Business_Rules.md",
            "Contact administrator for permissions & approvals",
        ],
    ),
]


def add_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title.replace("\n", " — ") if "\n" in title else title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)


def main() -> None:
    docs = Path(__file__).resolve().parent
    out = docs / "CA_Office_Suite_Employee_Presentation.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slide_layouts[0]
    first_title, first_bullets = SLIDES[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "CA Office Suite"
    slide.placeholders[1].text = "Business Rules & Logic\n" + "\n".join(first_bullets)

    for title, bullets in SLIDES[1:]:
        add_slide(prs, title, bullets)

    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
